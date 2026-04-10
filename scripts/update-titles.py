#!/usr/bin/env python3
"""Update slide titles in-place without touching anything else."""
from pptx import Presentation

path = "../out/piloc-motion-agent-presentation.pptx"
prs  = Presentation(path)

# { slide_index (0-based): (shape_index, new_title) }
UPDATES = {
    1: (4, "MP4 LinkedIn & YouTube, musique, voix off — tout généré en automatique."),
    2: (4, "De tes fichiers au MP4 publié, en 5 étapes."),
    3: (4, "Dépose tes fichiers dans in/ — l'agent analyse tout."),
    4: (4, "Décris ta feature dans le chat — l'agent structure la vidéo."),
    5: (4, "Dis à l'agent quelle animation utiliser — il l'intègre."),
}

for slide_idx, (shape_idx, new_text) in UPDATES.items():
    slide  = prs.slides[slide_idx]
    shape  = slide.shapes[shape_idx]
    tf     = shape.text_frame

    # Clear all paragraphs then restore one with the new text,
    # preserving the font properties of the first run.
    first_para = tf.paragraphs[0]
    first_run  = first_para.runs[0]

    # Capture existing formatting
    font_size  = first_run.font.size
    font_bold  = first_run.font.bold
    font_color = first_run.font.color.rgb
    font_name  = first_run.font.name
    align      = first_para.alignment

    # Replace text
    first_run.text = new_text

    # Restore formatting (pptx sometimes resets on text change)
    first_run.font.size       = font_size
    first_run.font.bold       = font_bold
    first_run.font.color.rgb  = font_color
    first_run.font.name       = font_name
    first_para.alignment      = align

    # Remove any extra runs or paragraphs that may exist
    for extra_run in first_para.runs[1:]:
        extra_run.text = ""
    for extra_para in tf.paragraphs[1:]:
        for run in extra_para.runs:
            run.text = ""

    print(f"  Slide {slide_idx+1} → {new_text[:60]}…")

prs.save(path)
print(f"\n✅  Saved in-place → {path}")
