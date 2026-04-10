#!/usr/bin/env python3
"""Read all text boxes per slide to find title locations."""
from pptx import Presentation

path = "../out/piloc-motion-agent-presentation.pptx"
prs  = Presentation(path)

for si, slide in enumerate(prs.slides):
    print(f"\n── Slide {si+1} ──────────────────────────")
    for shi, shape in enumerate(slide.shapes):
        if not shape.has_text_frame:
            continue
        full = " ".join(
            run.text
            for para in shape.text_frame.paragraphs
            for run in para.runs
        ).strip()
        if full:
            print(f"  [{shi}] size={shape.text_frame.paragraphs[0].runs[0].font.size} | {repr(full[:80])}")
