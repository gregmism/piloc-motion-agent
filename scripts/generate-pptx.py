#!/usr/bin/env python3
"""Generate Piloc Motion Agent presentation PPTX — v2."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ─── COLORS ────────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x0E, 0x10, 0x29)
MIST      = RGBColor(0x9D, 0xB8, 0xEE)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
WHITE72   = RGBColor(0xB8, 0xBD, 0xCE)
SUCCESS   = RGBColor(0x32, 0x95, 0x47)
CARD_BG   = RGBColor(0x16, 0x19, 0x38)
CHIP_BG   = RGBColor(0x1C, 0x21, 0x48)
ARROW_BG  = RGBColor(0x1A, 0x1E, 0x42)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ─── HELPERS ───────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    sh.line.fill.background()
    sh.line.width = 0
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_rgb
    return sh

def add_bg(slide):
    add_rect(slide, 0, 0, W, H, NAVY)
    # left accent bar
    add_rect(slide, 0, 0, Inches(0.07), H, MIST)

def add_text(slide, text, x, y, w, h,
             size=20, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False, spacing=None):
    tf = slide.shapes.add_textbox(x, y, w, h)
    tf.word_wrap = wrap
    frame = tf.text_frame
    if spacing:
        frame.paragraphs[0].space_before = Pt(spacing)
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Inter"
    return tf

def add_chip(slide, x, y, text, size=11):
    cw = Inches(len(text) * 0.115 + 0.55)
    ch = Inches(0.36)
    add_rect(slide, x, y, cw, ch, CHIP_BG)
    add_text(slide, text.upper(), x + Inches(0.18), y + Inches(0.06),
             cw - Inches(0.36), ch,
             size=size, bold=True, color=MIST)
    return cw

def add_card(slide, x, y, w, h, color=None):
    add_rect(slide, x, y, w, h, color or CARD_BG)

# Slide header helper (chip + title)
def add_header(slide, chip_text, title, title_color=WHITE):
    add_chip(slide, Inches(1.3), Inches(0.55), chip_text)
    add_text(slide, title,
             Inches(1.3), Inches(1.1), Inches(11), Inches(0.85),
             size=44, bold=True, color=title_color)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
add_bg(s1)

add_chip(s1, Inches(1.3), Inches(2.0), "PILOC · MOTION AGENT", size=12)

add_text(s1, "Génère ta vidéo marketing",
         Inches(1.3), Inches(2.6), Inches(11), Inches(1.15),
         size=58, bold=True, color=WHITE)
add_text(s1, "en quelques minutes.",
         Inches(1.3), Inches(3.7), Inches(11), Inches(1.1),
         size=58, bold=True, color=MIST)
add_text(s1,
         "Alimenté par tes fichiers HTML, screenshots et playlist — rendu en MP4.",
         Inches(1.3), Inches(4.95), Inches(9.5), Inches(0.65),
         size=20, color=WHITE72)
add_text(s1, "Motion Design Agent  ·  Remotion + TypeScript + AI",
         Inches(1.3), Inches(6.85), Inches(11), Inches(0.4),
         size=13, color=MIST, italic=True)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — CE QUE L'AGENT PRODUIT  (single wide card, no example output)
# ═══════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
add_bg(s2)
add_header(s2, "OUTPUTS", "Ce que l'agent produit")

outputs = [
    ("▶", "Horizontal  16:9",     "YouTube & TV Salon Pro  —  1920 × 1080 px"),
    ("◼", "Carré  1:1",           "LinkedIn autoplay  —  1080 × 1080 px"),
    ("🎵", "Musique beat-synced", "BPM auto-détecté, animations calées sur le rythme"),
    ("🎙", "Voix off ElevenLabs", "Sync frame-accurate par scène"),
    ("✦", "Animations multiples","Springs · Stagger · Typewriter · Curseur · KPI count-up"),
    ("📦", "Fichier de sortie",   "MP4 prêt à publier dans le dossier  out/"),
]

card_x = Inches(1.3)
card_y = Inches(2.2)
add_card(s2, card_x, card_y, Inches(10.8), Inches(4.6))

for i, (icon, label, sub) in enumerate(outputs):
    oy = card_y + Inches(0.45 + i * 0.67)
    add_text(s2, icon,  card_x + Inches(0.35), oy, Inches(0.45), Inches(0.45),
             size=18, color=MIST)
    add_text(s2, label, card_x + Inches(0.9),  oy, Inches(3.2), Inches(0.35),
             size=17, bold=True, color=WHITE)
    add_text(s2, sub,   card_x + Inches(0.9),  oy + Inches(0.34), Inches(9.0), Inches(0.28),
             size=14, color=WHITE72)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — FLOWCHART OVERVIEW
# ═══════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
add_bg(s3)
add_header(s3, "OVERVIEW", "Comment ça fonctionne ?")

# 5 flowchart nodes + arrows
nodes = [
    ("📁", "Dossier  in/",      "Screenshots, HTML\nPlaylist musique",    CARD_BG),
    ("💬", "Prompt",            "Feature, format,\nmusique, voix off",    CARD_BG),
    ("✅", "Validation",        "Storyboard approuvé\navant le code",      CARD_BG),
    ("🎬", "Vidéo générée",     "MP4 dans  out/\n16:9 + Carré auto",      RGBColor(0x10, 0x28, 0x18)),
    ("🔄", "Feedback",          "Itération par couches\ndans le chat",     CARD_BG),
]

total_nodes = len(nodes)
node_w = Inches(2.05)
node_h = Inches(3.2)
arrow_w = Inches(0.5)
total_w = total_nodes * node_w + (total_nodes - 1) * arrow_w
start_x = (W - total_w) / 2
node_y  = Inches(2.3)

for i, (icon, title, desc, bg) in enumerate(nodes):
    nx = start_x + i * (node_w + arrow_w)

    # node card
    add_card(s3, nx, node_y, node_w, node_h, bg)

    # icon
    add_text(s3, icon,
             nx, node_y + Inches(0.3),
             node_w, Inches(0.6),
             size=28, align=PP_ALIGN.CENTER)

    # title
    is_output = title.startswith("Vidéo")
    add_text(s3, title,
             nx + Inches(0.1), node_y + Inches(1.0),
             node_w - Inches(0.2), Inches(0.55),
             size=15, bold=True, color=SUCCESS if is_output else WHITE,
             align=PP_ALIGN.CENTER, wrap=True)

    # desc
    add_text(s3, desc,
             nx + Inches(0.15), node_y + Inches(1.65),
             node_w - Inches(0.3), Inches(1.3),
             size=13, color=WHITE72, align=PP_ALIGN.CENTER, wrap=True)

    # arrow (not after last node)
    if i < total_nodes - 1:
        ax = nx + node_w + Inches(0.05)
        ay = node_y + node_h / 2 - Inches(0.25)
        add_text(s3, "›",
                 ax, ay, arrow_w, Inches(0.5),
                 size=28, bold=True, color=MIST, align=PP_ALIGN.CENTER)

# bottom note
add_text(s3,
         "L'agent valide chaque étape avec toi — tu gardes le contrôle à chaque phase.",
         Inches(1.3), Inches(6.1), Inches(10.7), Inches(0.45),
         size=14, color=WHITE72, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SETUP : CE QUE TU DÉPOSES DANS in/
# ═══════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
add_bg(s4)
add_header(s4, "COMMENT DÉMARRER — 1/2", "Setup  :  le dossier  in/")

setup_cards = [
    {
        "icon":  "🖼",
        "title": "HTML / Screenshots",
        "chip":  "in/*.html  ·  in/*.png",
        "lines": [
            "Prototype HTML de la feature",
            "Captures d'écran de l'app",
            "",
            "L'agent extrait données,",
            "labels et parcours.",
            "",
            "Redesigné marketing —",
            "jamais affiché tel quel.",
        ],
    },
    {
        "icon":  "🎵",
        "title": "Playlist musique",
        "chip":  "in/music/*.mp3",
        "lines": [
            "Dépose tes tracks MP3/WAV.",
            "BPM analysé auto.",
            "",
            "Toutes les animations",
            "se calent sur le beat.",
            "",
            "Facultatif.",
            "",
        ],
    },
]

for j, d in enumerate(setup_cards):
    cx = Inches(1.3 + j * 6.0)
    cy = Inches(2.15)
    cw = Inches(5.5)
    ch = Inches(4.6)
    add_card(s4, cx, cy, cw, ch)

    add_text(s4, d["icon"] + "  " + d["title"],
             cx + Inches(0.4), cy + Inches(0.35),
             cw - Inches(0.8), Inches(0.5),
             size=20, bold=True, color=MIST)
    add_chip(s4, cx + Inches(0.4), cy + Inches(0.98), d["chip"], size=11)

    for k, line in enumerate(d["lines"]):
        col = MIST if line.startswith("→") or line in ("Facultatif.", "jamais affiché tel quel.") else WHITE if line else WHITE72
        add_text(s4, line,
                 cx + Inches(0.4), cy + Inches(1.55 + k * 0.37),
                 cw - Inches(0.8), Inches(0.38),
                 size=14, color=col)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — CHAT : LE PROMPT
# ═══════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
add_bg(s5)
add_header(s5, "COMMENT DÉMARRER — 2/2", "Chat  :  ton prompt")

# Left: what to say
left_x = Inches(1.3)
left_y = Inches(2.15)
add_card(s5, left_x, left_y, Inches(5.2), Inches(4.6))

add_text(s5, "Ce que tu précises",
         left_x + Inches(0.4), left_y + Inches(0.35),
         Inches(4.5), Inches(0.45), size=18, bold=True, color=MIST)

prompt_items = [
    ("Feature à promouvoir ?",         "ex. Campagnes — parcours de création"),
    ("Valeur client à montrer ?",       "ex. Gagner du temps sur les relances"),
    ("Musique : oui / non + titre ?",   "L'agent liste tes tracks disponibles"),
    ("Voix off : oui / non ?",          "ElevenLabs, sync frame-accurate"),
    ("Format : horizontal ou carré ?",  "16:9 YouTube / 1:1 LinkedIn"),
]
for k, (q, hint) in enumerate(prompt_items):
    oy = left_y + Inches(1.05 + k * 0.66)
    add_text(s5, q,    left_x + Inches(0.4), oy, Inches(4.4), Inches(0.32),
             size=15, bold=True, color=WHITE)
    add_text(s5, hint, left_x + Inches(0.4), oy + Inches(0.3), Inches(4.4), Inches(0.28),
             size=12, color=WHITE72)

# Right: example exchange
rx = Inches(7.1)
ry = Inches(2.15)
add_card(s5, rx, ry, Inches(5.0), Inches(2.05))
add_text(s5, "Vous",
         rx + Inches(0.35), ry + Inches(0.2),
         Inches(1.5), Inches(0.32), size=12, bold=True, color=MIST)
user_msg = (
    "Je veux promouvoir la feature Campagnes,\n"
    "le parcours de création. Expliquer la valeur client.\n"
    "Musique : oui → « Midnight Drive »\n"
    "Voix off : non · Format : horizontal + carré"
)
add_text(s5, user_msg,
         rx + Inches(0.35), ry + Inches(0.6), Inches(4.25), Inches(1.3),
         size=13, color=WHITE, wrap=True)

# agent reply
ary = ry + Inches(2.25)
add_card(s5, rx, ary, Inches(5.0), Inches(2.35))
add_text(s5, "Agent",
         rx + Inches(0.35), ary + Inches(0.2),
         Inches(1.5), Inches(0.32), size=12, bold=True, color=SUCCESS)
agent_msg = (
    "BPM détecté : 115 — Beat = 16f · Bar = 63f\n\n"
    "Story Map extraite. Voici le storyboard\n"
    "en 5 scènes — dis-moi si ça te convient ✅"
)
add_text(s5, agent_msg,
         rx + Inches(0.35), ary + Inches(0.6), Inches(4.25), Inches(1.5),
         size=13, color=WHITE, wrap=True)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — BIBLIOTHÈQUE D'ANIMATIONS
# ═══════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
add_bg(s6)
add_header(s6, "ANIMATIONS — BIBLIOTHÈQUE", "Précise ce que tu veux animer")

add_text(s6,
         "L'agent anime par défaut. Indique dans ton prompt les effets supplémentaires souhaités.",
         Inches(1.3), Inches(2.05), Inches(10.8), Inches(0.42),
         size=16, color=WHITE72)

anims = [
    # (name, desc, usage, col)
    ("Typewriter",     "Le texte s'écrit lettre par lettre",       "Champs de formulaire, emails",   0),
    ("Curseur souris", "Reproduit déplacements et clics",           "Montrer l'UX de façon naturelle",0),
    ("KPI Count-up",   "Les chiffres s'incrémentent",               "Métriques, dashboards",          0),
    ("Spring Pop",     "Apparition avec rebond maîtrisé",           "Cartes, modales, badges",        0),
    ("Progress Bar",   "Barre qui se remplit progressivement",      "Taux, pipeline, progression",    0),
    ("Stagger",        "Éléments qui apparaissent en cascade",      "Lignes de tableau, pills",       1),
    ("Badge Flip",     "Statut qui change — rouge → vert",          "Beat 3, moment clé",             1),
    ("Highlight Row",  "Ligne de tableau mise en évidence",         "Sélection, alerte",              1),
    ("SVG Chart",      "Courbe qui se trace en direct",             "Données financières, KPIs",      1),
    ("Beat-sync",      "Animations calées sur le BPM de la musique","Transitions, entrées de scène",  1),
]

col_x = [Inches(1.3), Inches(7.2)]
for i, (name, desc, usage, col) in enumerate(anims):
    row = i % 5
    ox  = col_x[col]
    oy  = Inches(2.7 + row * 0.88)

    # left mist bar
    add_rect(s6, ox, oy + Inches(0.06), Inches(0.05), Inches(0.6), MIST)
    add_text(s6, name, ox + Inches(0.22), oy, Inches(5.5), Inches(0.36),
             size=16, bold=True, color=WHITE)
    add_text(s6, desc, ox + Inches(0.22), oy + Inches(0.34), Inches(5.5), Inches(0.28),
             size=13, color=WHITE72)
    add_text(s6, "→ " + usage, ox + Inches(0.22), oy + Inches(0.6), Inches(5.5), Inches(0.24),
             size=12, color=MIST)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — WORKFLOW EN 7 PHASES
# ═══════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
add_bg(s7)
add_header(s7, "WORKFLOW — 7 PHASES", "De l'idée au MP4 publié")

phases = [
    ("0",   "Intake",      "Musique, scénario,\nformat",     False),
    ("1",   "Analyze",     "Story Map depuis\ntes fichiers", False),
    ("2",   "Storyboard",  "Arc 4 beats\n+ timeline",        False),
    ("3",   "Approval",    "Validation\navant le code",       False),
    ("3.5", "Pre-plan",    "Math, composants\nbeat grid",     False),
    ("4",   "Execute",     "TSX + render\nMP4",               True),
    ("5",   "Review",      "Feedback couche\npar couche",     False),
]

n = len(phases)
pw  = Inches(1.5)
ph  = Inches(3.8)
gap = Inches(0.3)
total_row = n * pw + (n - 1) * gap
sx = (W - total_row) / 2
py = Inches(2.4)

for i, (num, name, desc, is_code) in enumerate(phases):
    px = sx + i * (pw + gap)
    bg = RGBColor(0x10, 0x28, 0x18) if is_code else CARD_BG
    add_card(s7, px, py, pw, ph, bg)

    # number badge
    badge_col = SUCCESS if is_code else MIST
    add_rect(s7, px + Inches(0.55), py + Inches(0.22), Inches(0.4), Inches(0.4), badge_col)
    add_text(s7, num,
             px + Inches(0.55), py + Inches(0.2),
             Inches(0.4), Inches(0.42),
             size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    add_text(s7, name,
             px + Inches(0.1), py + Inches(0.82),
             pw - Inches(0.2), Inches(0.42),
             size=14, bold=True,
             color=SUCCESS if is_code else WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s7, desc,
             px + Inches(0.1), py + Inches(1.35),
             pw - Inches(0.2), Inches(1.6),
             size=12, color=WHITE72,
             align=PP_ALIGN.CENTER, wrap=True)

    if i < n - 1:
        ax = px + pw + Inches(0.02)
        ay = py + ph / 2 - Inches(0.22)
        add_text(s7, "›", ax, ay, gap - Inches(0.02), Inches(0.44),
                 size=20, bold=True, color=MIST, align=PP_ALIGN.CENTER)

add_text(s7,
         "Chaque ✅ déclenche la phase suivante — tu peux réviser n'importe quelle couche à tout moment.",
         Inches(1.3), Inches(6.55), Inches(10.7), Inches(0.45),
         size=14, color=WHITE72, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — CTA
# ═══════════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK)
add_bg(s8)
add_rect(s8, 0, 0, Inches(0.07), H, MIST)

add_chip(s8, Inches(1.3), Inches(1.85), "PILOC · MOTION AGENT", size=12)
add_text(s8, "Prêt à créer ta première",
         Inches(1.3), Inches(2.45), Inches(11), Inches(1.1),
         size=58, bold=True, color=WHITE)
add_text(s8, "vidéo marketing ?",
         Inches(1.3), Inches(3.5), Inches(11), Inches(1.1),
         size=58, bold=True, color=MIST)

steps = [
    ("1", "Dépose ton prototype HTML ou screenshots dans  in/"),
    ("2", "Ajoute une musique dans  in/music/  (facultatif)"),
    ("3", "Lance l'agent — décris ta feature dans le chat"),
    ("4", "Récupère ton  .mp4  dans  out/  — prêt à publier"),
]
for i, (num, text) in enumerate(steps):
    oy = Inches(4.85 + i * 0.47)
    add_rect(s8, Inches(1.3), oy + Inches(0.05), Inches(0.28), Inches(0.28),
             MIST if i < 3 else SUCCESS)
    add_text(s8, num, Inches(1.3), oy + Inches(0.03),
             Inches(0.28), Inches(0.3),
             size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s8, text, Inches(1.75), oy,
             Inches(9.5), Inches(0.38),
             size=16, color=WHITE if i < 3 else MIST, bold=(i == 3))

# ─── SAVE ──────────────────────────────────────────────────────────────────
out_path = "../out/piloc-motion-agent-presentation.pptx"
os.makedirs("../out", exist_ok=True)
prs.save(out_path)
print(f"✅  Saved → {out_path}")
